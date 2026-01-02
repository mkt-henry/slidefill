#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 템플릿 변환 스크립트
엑셀 파일의 'Start Word'와 'End Word' 열을 읽어 PPT 템플릿의 텍스트를 치환하고,
이미지를 지정된 자리표시자에 삽입하는 기능을 제공합니다.

사용법:
    python convert_ppt.py <template_path> <excel_path> <output_path> [image_mappings_json]
"""

import sys
import os

# 타입 체킹 비활성화
os.environ['PYTHONDONTWRITEBYTECODE'] = '1'

import json
import pandas as pd

try:
    from pptx import Presentation
except ImportError as e:
    print(f"Error: python-pptx 패키지를 찾을 수 없습니다: {e}", file=sys.stderr)
    sys.exit(1)
from pptx.util import Inches
from PIL import Image as PILImage
import io
import os


def parse_excel(excel_path):
    """
    엑셀 파일에서 'Start Word'와 'End Word' 열을 읽어 치환 딕셔너리를 생성합니다.
    
    Args:
        excel_path: 엑셀 파일 경로
        
    Returns:
        dict: {start_word: end_word} 형태의 딕셔너리
        int: 치환 단어 쌍의 개수
    """
    try:
        df = pd.read_excel(excel_path, engine='openpyxl')
        
        # 열 이름 정규화 (대소문자 및 공백 무시)
        df.columns = df.columns.str.strip().str.lower()
        
        # 'start word'와 'end word' 열 찾기
        start_col = None
        end_col = None
        
        for col in df.columns:
            if 'start' in col and 'word' in col:
                start_col = col
            if 'end' in col and 'word' in col:
                end_col = col
        
        if start_col is None or end_col is None:
            raise ValueError("엑셀 파일에 'Start Word'와 'End Word' 열이 필요합니다.")
        
        # NaN 값 제거 및 문자열 변환
        replacements = {}
        for _, row in df.iterrows():
            start = str(row[start_col]).strip() if pd.notna(row[start_col]) else ""
            end = str(row[end_col]).strip() if pd.notna(row[end_col]) else ""
            if start and end:
                replacements[start] = end
        
        return replacements, len(replacements)
    
    except Exception as e:
        raise Exception(f"엑셀 파일 파싱 중 오류 발생: {str(e)}")


def replace_text_in_shape(shape, replacements):
    """
    도형 내의 텍스트를 치환합니다.
    
    Args:
        shape: PPT 도형 객체
        replacements: 치환 딕셔너리
    """
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)


def replace_text_in_table(table, replacements):
    """
    테이블 내의 텍스트를 치환합니다.
    
    Args:
        table: PPT 테이블 객체
        replacements: 치환 딕셔너리
    """
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)


def insert_image_to_placeholder(slide, placeholder_name, image_path):
    """
    슬라이드의 특정 자리표시자에 이미지를 삽입합니다.
    
    Args:
        slide: PPT 슬라이드 객체
        placeholder_name: 자리표시자 이름 (예: {{profile_image}})
        image_path: 삽입할 이미지 파일 경로
    """
    for shape in slide.shapes:
        # 텍스트 프레임이 있는 도형에서 자리표시자 이름 찾기
        if shape.has_text_frame:
            if placeholder_name in shape.text_frame.text:
                # 도형의 위치와 크기 저장
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # 기존 도형 삭제
                sp = shape.element
                sp.getparent().remove(sp)
                
                # 이미지 삽입
                try:
                    pic = slide.shapes.add_picture(image_path, left, top, width, height)
                    return True
                except Exception as e:
                    print(f"이미지 삽입 중 오류: {str(e)}")
                    return False
    
    return False


def convert_ppt(template_path, excel_path, output_path, image_mappings=None):
    """
    PPT 템플릿을 변환합니다.
    
    Args:
        template_path: 템플릿 PPT 파일 경로
        excel_path: 엑셀 데이터 파일 경로
        output_path: 출력 PPT 파일 경로
        image_mappings: 이미지 매핑 정보 (dict)
                       예: {"{{profile_image}}": "/path/to/image.jpg"}
    
    Returns:
        dict: 변환 결과 정보
    """
    try:
        # 엑셀 파일 파싱
        replacements, word_pair_count = parse_excel(excel_path)
        
        if not replacements:
            raise ValueError("엑셀 파일에 유효한 치환 데이터가 없습니다.")
        
        # PPT 파일 로드
        prs = Presentation(template_path)
        
        # 슬라이드 수 확인
        slide_count = len(prs.slides)
        
        # 모든 슬라이드 순회하며 텍스트 치환
        for slide in prs.slides:
            for shape in slide.shapes:
                # 일반 도형의 텍스트 치환
                replace_text_in_shape(shape, replacements)
                
                # 테이블의 텍스트 치환
                if shape.has_table:
                    replace_text_in_table(shape.table, replacements)
                
                # 그룹 도형 내부 처리
                if shape.shape_type == 6:  # GROUP
                    for sub_shape in shape.shapes:
                        replace_text_in_shape(sub_shape, replacements)
        
        # 이미지 삽입 (있는 경우)
        images_inserted = 0
        if image_mappings:
            for placeholder_name, image_path in image_mappings.items():
                if os.path.exists(image_path):
                    for slide in prs.slides:
                        if insert_image_to_placeholder(slide, placeholder_name, image_path):
                            images_inserted += 1
        
        # 변환된 PPT 저장
        prs.save(output_path)
        
        return {
            "success": True,
            "slide_count": slide_count,
            "word_pair_count": word_pair_count,
            "images_inserted": images_inserted,
            "output_path": output_path
        }
    
    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }


def main():
    """메인 함수"""
    if len(sys.argv) < 4:
        print("사용법: python convert_ppt.py <template_path> <excel_path> <output_path> [image_mappings_json]")
        sys.exit(1)
    
    template_path = sys.argv[1]
    excel_path = sys.argv[2]
    output_path = sys.argv[3]
    image_mappings = None
    
    if len(sys.argv) >= 5:
        try:
            image_mappings = json.loads(sys.argv[4])
        except json.JSONDecodeError:
            print("이미지 매핑 JSON 파싱 오류")
            sys.exit(1)
    
    # 파일 존재 확인
    if not os.path.exists(template_path):
        print(f"템플릿 파일을 찾을 수 없습니다: {template_path}")
        sys.exit(1)
    
    if not os.path.exists(excel_path):
        print(f"엑셀 파일을 찾을 수 없습니다: {excel_path}")
        sys.exit(1)
    
    # 변환 실행
    result = convert_ppt(template_path, excel_path, output_path, image_mappings)
    
    # 결과 출력 (JSON 형태)
    print(json.dumps(result, ensure_ascii=False))
    
    if result["success"]:
        sys.exit(0)
    else:
        sys.exit(1)


if __name__ == "__main__":
    main()
